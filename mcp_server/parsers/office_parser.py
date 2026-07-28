"""
╭─╴ KNOWLEDGE-RAG OFFICE PARSER ╶────────────────────────────────╮
│                                                                │
│   Microsoft Office readers: DOCX, XLSX, PPTX.                  │
│                                                                │
╰────────────────────────────────────────────────────────────────╯

    ┌─ Author  ·  Ailton Rocha (Lyon.)
    └─ Date    ·  2026-07-27

Each format is its own parser class because the underlying libraries
(``python-docx``, ``openpyxl``, ``python-pptx``) are independent and can
be missing individually. Bundling them in a single module keeps the
parsers/ directory tidy without coupling the availability of one
optional dep to another — if only ``python-docx`` is installed, DOCX
files parse and XLSX/PPTX raise :class:`ImportError` at parse time with
the exact ``pip`` command.

DOCX headings are converted to Markdown-style ``#`` prefixes so
retrieval sees a consistent header shape whether the source was ``.md``
or ``.docx``. XLSX and PPTX get ``## Sheet: …`` / ``## Slide N`` section
markers for the same reason.
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Any, Dict

from .base import ParserResult
from .registry import register_parser

__all__ = ["DocxParser", "XlsxParser", "PptxParser"]


try:
    import docx  # type: ignore[import-not-found]  # python-docx

    _HAS_DOCX = True
except ImportError:  # pragma: no cover
    _HAS_DOCX = False

try:
    import openpyxl  # type: ignore[import-not-found]

    _HAS_XLSX = True
except ImportError:  # pragma: no cover
    _HAS_XLSX = False

try:
    from pptx import Presentation  # type: ignore[import-not-found]

    _HAS_PPTX = True
except ImportError:  # pragma: no cover
    _HAS_PPTX = False


@register_parser
class DocxParser:
    """Parse ``.docx`` files, preserving heading hierarchy and tables."""

    extensions: frozenset[str] = frozenset({".docx"})
    display_name: str = "docx"
    optional_deps: tuple[str, ...] = ("python-docx",)

    def can_parse(self, path: Path) -> bool:
        """Return ``True`` for ``.docx`` files.

        Args:
            path: Candidate file path.

        Returns:
            bool: Membership test against :attr:`extensions`.
        """
        return path.suffix.lower() in self.extensions

    def parse(self, path: Path) -> ParserResult:
        """Extract paragraphs and tables from a Word document.

        Headings become Markdown ``#`` lines whose depth matches the
        Word heading level. Tables are flattened row-by-row with `` | ``
        separators.

        Args:
            path: Existing ``.docx`` file.

        Returns:
            ParserResult: ``(content, metadata)`` where metadata carries
            ``type``, ``title``, ``file_size``, ``modified``,
            ``paragraphs`` count, ``tables`` count.

        Raises:
            ImportError: When ``python-docx`` is not installed.
        """
        if not _HAS_DOCX:
            raise ImportError("python-docx not installed. Install with: pip install python-docx")

        doc = docx.Document(path)
        metadata: Dict[str, Any] = {
            "type": "docx",
            "title": path.stem,
            "file_size": path.stat().st_size,
            "modified": datetime.fromtimestamp(path.stat().st_mtime).isoformat(),
            "paragraphs": len(doc.paragraphs),
            "tables": len(doc.tables),
        }

        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Preserve heading structure as markdown
                if para.style and para.style.name.startswith("Heading"):
                    try:
                        level = int(para.style.name.split()[-1])
                        parts.append(f"{'#' * level} {text}")
                    except (ValueError, IndexError):
                        parts.append(f"## {text}")
                else:
                    parts.append(text)

        # Extract tables as markdown
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                parts.append("\n".join(rows))

        content = "\n\n".join(parts)
        return content, metadata


@register_parser
class XlsxParser:
    """Parse ``.xlsx`` workbooks sheet-by-sheet."""

    extensions: frozenset[str] = frozenset({".xlsx"})
    display_name: str = "xlsx"
    optional_deps: tuple[str, ...] = ("openpyxl",)

    def can_parse(self, path: Path) -> bool:
        """Return ``True`` for ``.xlsx`` files.

        Args:
            path: Candidate file path.

        Returns:
            bool: Membership test against :attr:`extensions`.
        """
        return path.suffix.lower() in self.extensions

    def parse(self, path: Path) -> ParserResult:
        """Flatten each sheet into a ``## Sheet: name`` section.

        Runs ``openpyxl`` in ``read_only=True, data_only=True`` mode so
        formulas resolve to their cached values and memory stays bounded
        on multi-megabyte workbooks. Rows that reduce to only pipe
        separators (empty cells) are dropped to keep the indexed text
        signal-heavy.

        Args:
            path: Existing ``.xlsx`` file.

        Returns:
            ParserResult: ``(content, metadata)`` where metadata carries
            ``type``, ``title``, ``file_size``, ``modified``, and
            ``sheets`` (list of sheet names).

        Raises:
            ImportError: When ``openpyxl`` is not installed.
        """
        if not _HAS_XLSX:
            raise ImportError("openpyxl not installed. Install with: pip install openpyxl")

        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        metadata: Dict[str, Any] = {
            "type": "xlsx",
            "title": path.stem,
            "file_size": path.stat().st_size,
            "modified": datetime.fromtimestamp(path.stat().st_mtime).isoformat(),
            "sheets": wb.sheetnames,
        }

        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## Sheet: {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                line = " | ".join(cells).strip()
                if line and line != " | " * (len(cells) - 1):
                    parts.append(line)

        wb.close()
        content = "\n\n".join(parts)
        return content, metadata


@register_parser
class PptxParser:
    """Parse ``.pptx`` decks slide-by-slide."""

    extensions: frozenset[str] = frozenset({".pptx"})
    display_name: str = "pptx"
    optional_deps: tuple[str, ...] = ("python-pptx",)

    def can_parse(self, path: Path) -> bool:
        """Return ``True`` for ``.pptx`` files.

        Args:
            path: Candidate file path.

        Returns:
            bool: Membership test against :attr:`extensions`.
        """
        return path.suffix.lower() in self.extensions

    def parse(self, path: Path) -> ParserResult:
        """Extract text from every shape with a text frame.

        Args:
            path: Existing ``.pptx`` file.

        Returns:
            ParserResult: ``(content, metadata)`` where content is each
            slide's text joined by blank lines with ``## Slide N``
            markers, and metadata carries ``type``, ``title``,
            ``file_size``, ``modified``, and ``slides`` count.

        Raises:
            ImportError: When ``python-pptx`` is not installed.
        """
        if not _HAS_PPTX:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")

        prs = Presentation(path)
        metadata: Dict[str, Any] = {
            "type": "pptx",
            "title": path.stem,
            "file_size": path.stat().st_size,
            "modified": datetime.fromtimestamp(path.stat().st_mtime).isoformat(),
            "slides": len(prs.slides),
        }

        parts = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                parts.append(f"## Slide {i + 1}\n" + "\n".join(slide_texts))

        content = "\n\n".join(parts)
        return content, metadata
